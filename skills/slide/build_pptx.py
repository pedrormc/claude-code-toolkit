"""
build_pptx.py — Exporta apresentacao Singular pra .pptx (PowerPoint).

Uso:
    python build_pptx.py content.json saida.pptx

Le o mesmo content.json do build_html.py e gera um .pptx editavel com:
- Paleta Singular (cobre #E64E10, off-white #F7EEEB, preto profundo #1C1C1C)
- Fonte Urbanist (se instalada no sistema; fallback Calibri)
- Logo Singular na capa
- Layout 16:9, fundo preto profundo

Renderizacao simplificada — o HTML continua sendo a fonte de verdade visual.
PPTX serve pra clientes/contextos onde precisa formato editavel.

Dependencia: python-pptx (pip install python-pptx)
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Erro: python-pptx nao instalado. Rode: pip install python-pptx", file=sys.stderr)
    sys.exit(3)

# Paleta Singular
INK = RGBColor(0x1C, 0x1C, 0x1C)
INK_SOFT = RGBColor(0x2A, 0x27, 0x27)
PAPER = RGBColor(0xF7, 0xEE, 0xEB)
COPPER = RGBColor(0xE6, 0x4E, 0x10)
STEEL = RGBColor(0x5B, 0x4B, 0x48)
MUTED = RGBColor(0xB7, 0xAF, 0xAB)

FONT = "Urbanist"
FONT_FALLBACK = "Calibri"

ASSETS_DIR = Path(__file__).parent / "assets"
LOGO_DARK = ASSETS_DIR / "isotipo-singular-dark.png"


def set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=PAPER, align=PP_ALIGN.LEFT, italic=False, font=FONT, letter_spacing=0):
    # letter_spacing nao tem suporte direto em python-pptx — parametro mantido pra compat
    _ = letter_spacing
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_rect(slide, left, top, width, height, fill_color, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.5)
    return shape


def add_eyebrow(slide, left, top, text):
    add_text(slide, left, top, Inches(8), Inches(0.3),
             text.upper(), size=10, bold=True, color=COPPER, letter_spacing=2)


def render_cover(slide, data, total):
    set_slide_bg(slide, INK)
    add_eyebrow(slide, Inches(0.6), Inches(0.6), data.get("eyebrow", "Apresentacao"))

    add_text(slide, Inches(0.6), Inches(2.0), Inches(12), Inches(3),
             data.get("hero", "Singular"), size=72, bold=True, color=PAPER)

    add_text(slide, Inches(0.6), Inches(4.2), Inches(12), Inches(0.7),
             data.get("subtitle", ""), size=20, color=MUTED)

    add_rect(slide, Inches(0.6), Inches(6.2), Inches(12), Pt(0.5), STEEL)
    add_text(slide, Inches(0.6), Inches(6.4), Inches(8), Inches(0.4),
             data.get("footer_left", ""), size=10, color=MUTED)
    add_text(slide, Inches(8.6), Inches(6.4), Inches(4), Inches(0.4),
             data.get("footer_right", ""), size=10, color=COPPER, align=PP_ALIGN.RIGHT, bold=True)

    if LOGO_DARK.exists():
        slide.shapes.add_picture(str(LOGO_DARK), Inches(11.2), Inches(0.5), height=Inches(0.6))


def render_summary_cards(slide, data, total, idx):
    set_slide_bg(slide, INK)
    add_slide_number(slide, idx, total)
    add_eyebrow(slide, Inches(0.6), Inches(0.6), data.get("eyebrow", ""))
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(1.3),
             data.get("title", ""), size=32, bold=True, color=PAPER)

    cards = data.get("cards", [])
    n = max(1, len(cards))
    avail = Inches(12.4)
    gap = Inches(0.25)
    card_w = Emu(int((avail - gap * (n - 1)) / n)) if n > 1 else avail
    card_h = Inches(3.4)
    top = Inches(2.7)
    left0 = Inches(0.6)

    for i, c in enumerate(cards):
        left = Emu(int(left0) + i * (int(card_w) + int(gap)))
        bg = INK_SOFT
        rect = add_rect(slide, left, top, card_w, card_h, bg, line=STEEL)
        if c.get("recommended"):
            rect.line.color.rgb = COPPER
            rect.line.width = Pt(1.5)
        add_text(slide, left + Inches(0.25), top + Inches(0.25), card_w - Inches(0.5), Inches(0.4),
                 c.get("tag", "").upper(), size=11, bold=True, color=COPPER)
        add_text(slide, left + Inches(0.25), top + Inches(0.7), card_w - Inches(0.5), Inches(0.9),
                 c.get("title", ""), size=18, bold=True, color=PAPER)
        add_text(slide, left + Inches(0.25), top + Inches(1.7), card_w - Inches(0.5), card_h - Inches(2),
                 c.get("body", ""), size=12, color=MUTED)


def render_section_detail(slide, data, total, idx):
    set_slide_bg(slide, INK)
    add_slide_number(slide, idx, total)

    badge = data.get("badge", "")
    if badge:
        bw = Inches(1.4)
        add_rect(slide, Inches(0.6), Inches(0.6), bw, Inches(0.4), COPPER)
        add_text(slide, Inches(0.6), Inches(0.62), bw, Inches(0.4),
                 badge.upper(), size=10, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(1.15), Inches(12), Inches(1.2),
             data.get("title", ""), size=30, bold=True, color=PAPER)
    add_text(slide, Inches(0.6), Inches(2.2), Inches(12), Inches(0.5),
             data.get("subtitle", ""), size=14, color=MUTED)

    cards = data.get("cards", [])
    n = max(1, len(cards))
    avail = Inches(12.4)
    gap = Inches(0.25)
    card_w = Emu(int((avail - gap * (n - 1)) / n)) if n > 1 else avail
    card_h = Inches(2.6)
    top = Inches(3.0)
    left0 = Inches(0.6)

    for i, c in enumerate(cards):
        left = Emu(int(left0) + i * (int(card_w) + int(gap)))
        add_rect(slide, left, top, card_w, card_h, INK_SOFT, line=STEEL)
        add_text(slide, left + Inches(0.25), top + Inches(0.2), card_w - Inches(0.5), Inches(0.4),
                 c.get("label", "").upper(), size=10, bold=True, color=COPPER)
        add_text(slide, left + Inches(0.25), top + Inches(0.7), card_w - Inches(0.5), card_h - Inches(0.9),
                 c.get("value", ""), size=13, color=PAPER)

    if data.get("highlight_box"):
        hb_top = Inches(5.8)
        add_rect(slide, Inches(0.6), hb_top, Inches(12.4), Inches(1), INK_SOFT, line=COPPER)
        add_text(slide, Inches(0.85), hb_top + Inches(0.2), Inches(12), Inches(0.7),
                 data["highlight_box"], size=16, bold=True, color=PAPER)


def render_compare_table(slide, data, total, idx):
    set_slide_bg(slide, INK)
    add_slide_number(slide, idx, total)
    add_eyebrow(slide, Inches(0.6), Inches(0.6), data.get("eyebrow", ""))
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.9),
             data.get("title", ""), size=28, bold=True, color=PAPER)

    cols = data.get("columns", [])
    rows = data.get("rows", [])
    base_idx = data.get("base_col_index", 1)
    if not cols or not rows:
        return

    table_left = Inches(0.6)
    table_top = Inches(2.2)
    table_w = Inches(12.4)
    n_cols = len(cols)
    n_rows = len(rows) + 1
    table_h = Emu(int(Inches(0.55)) * n_rows)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_w, table_h)
    tbl = tbl_shape.table

    for j, col in enumerate(cols):
        cell = tbl.cell(0, j)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = INK_SOFT if j != base_idx else COPPER
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = str(col)
        run.font.name = FONT
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = INK if j == base_idx else MUTED

    for i, row in enumerate(rows):
        for j, raw_cell in enumerate(row):
            value = raw_cell.get("value", "") if isinstance(raw_cell, dict) else raw_cell
            cell = tbl.cell(i + 1, j)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                RGBColor(0x29, 0x1B, 0x18) if j == base_idx else INK
            )
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(value)
            run.font.name = FONT
            run.font.size = Pt(11)
            run.font.bold = (j == base_idx) or (isinstance(raw_cell, dict) and raw_cell.get("strong"))
            run.font.color.rgb = PAPER if j != 0 else MUTED


def render_text_section(slide, data, total, idx):
    set_slide_bg(slide, INK)
    add_slide_number(slide, idx, total)
    add_eyebrow(slide, Inches(0.6), Inches(0.6), data.get("eyebrow", ""))
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(1),
             data.get("title", ""), size=28, bold=True, color=PAPER)
    if data.get("subtitle"):
        add_text(slide, Inches(0.6), Inches(2.0), Inches(12), Inches(0.5),
                 data["subtitle"], size=14, color=MUTED)

    y = Inches(2.7)
    for p in data.get("paragraphs", []):
        add_text(slide, Inches(0.6), y, Inches(12.4), Inches(0.6),
                 p, size=14, color=PAPER)
        y = Emu(int(y) + int(Inches(0.7)))
    for b in data.get("bullets", []):
        add_text(slide, Inches(0.6), y, Inches(12.4), Inches(0.5),
                 "—  " + b, size=13, color=PAPER)
        y = Emu(int(y) + int(Inches(0.5)))
    if data.get("highlight_box"):
        add_rect(slide, Inches(0.6), y + Inches(0.2), Inches(12.4), Inches(0.9),
                 INK_SOFT, line=COPPER)
        add_text(slide, Inches(0.85), y + Inches(0.4), Inches(12), Inches(0.6),
                 data["highlight_box"], size=15, bold=True, color=PAPER)


def render_recommendation(slide, data, total, idx):
    set_slide_bg(slide, INK)
    add_slide_number(slide, idx, total)
    add_text(slide, Inches(0.6), Inches(0.8), Inches(12.4), Inches(0.4),
             data.get("eyebrow", "Nossa recomendacao").upper(),
             size=11, bold=True, color=COPPER, align=PP_ALIGN.CENTER, letter_spacing=3)

    headline = data.get("headline", "")
    accent = data.get("accent", "")
    full = f"{headline} {accent}" if accent else headline
    add_text(slide, Inches(0.6), Inches(1.4), Inches(12.4), Inches(1.5),
             full, size=64, bold=True, color=PAPER, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.6), Inches(3.1), Inches(12.4), Inches(0.6),
             data.get("sub", ""), size=16, color=MUTED, align=PP_ALIGN.CENTER)

    items = data.get("items", [])
    grid_top = Inches(4.0)
    grid_left = Inches(1.5)
    grid_w = Inches(10.5)
    cell_w = Emu(int(grid_w) // 2 - int(Inches(0.15)))
    cell_h = Inches(1.3)
    for i, it in enumerate(items[:4]):
        col = i % 2
        row = i // 2
        x = Emu(int(grid_left) + col * (int(cell_w) + int(Inches(0.3))))
        y = Emu(int(grid_top) + row * (int(cell_h) + int(Inches(0.2))))
        add_rect(slide, x, y, cell_w, cell_h, INK_SOFT)
        add_rect(slide, x, y, Inches(0.05), cell_h, COPPER)
        add_text(slide, x + Inches(0.15), y + Inches(0.1), Inches(0.7), Inches(0.5),
                 f"{i+1:02d}", size=20, bold=True, color=COPPER)
        add_text(slide, x + Inches(0.85), y + Inches(0.1), cell_w - Inches(1), Inches(0.4),
                 it.get("title", ""), size=12, bold=True, color=PAPER)
        add_text(slide, x + Inches(0.85), y + Inches(0.55), cell_w - Inches(1), cell_h - Inches(0.7),
                 it.get("body", ""), size=10, color=MUTED)


def render_next_steps(slide, data, total, idx):
    set_slide_bg(slide, INK)
    add_slide_number(slide, idx, total)
    add_eyebrow(slide, Inches(0.6), Inches(0.6), data.get("eyebrow", "Proximos passos"))
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(1),
             data.get("title", ""), size=28, bold=True, color=PAPER)

    cards = data.get("cards", [])
    n = max(1, len(cards))
    avail = Inches(12.4)
    gap = Inches(0.25)
    card_w = Emu(int((avail - gap * (n - 1)) / n)) if n > 1 else avail
    card_h = Inches(3.0)
    top = Inches(2.7)
    left0 = Inches(0.6)

    for i, c in enumerate(cards):
        left = Emu(int(left0) + i * (int(card_w) + int(gap)))
        add_rect(slide, left, top, card_w, card_h, INK_SOFT, line=STEEL)
        add_text(slide, left + Inches(0.25), top + Inches(0.25), card_w - Inches(0.5), Inches(0.7),
                 f"{i+1:02d}", size=28, bold=True, color=COPPER)
        add_text(slide, left + Inches(0.25), top + Inches(1.0), card_w - Inches(0.5), Inches(0.5),
                 c.get("head", "").upper(), size=14, bold=True, color=PAPER, letter_spacing=1)
        add_text(slide, left + Inches(0.25), top + Inches(1.5), card_w - Inches(0.5), Inches(1.2),
                 c.get("question", ""), size=12, color=MUTED)
        if c.get("options"):
            add_text(slide, left + Inches(0.25), top + card_h - Inches(0.45), card_w - Inches(0.5), Inches(0.4),
                     c["options"], size=10, color=COPPER, bold=True)

    if data.get("cta"):
        cta_top = Inches(6.1)
        add_rect(slide, Inches(2), cta_top, Inches(9.6), Inches(0.8), COPPER)
        add_text(slide, Inches(2), cta_top + Inches(0.2), Inches(9.6), Inches(0.5),
                 data["cta"], size=16, bold=True, color=INK, align=PP_ALIGN.CENTER)


def render_closing(slide, data, total, idx):
    set_slide_bg(slide, INK)
    add_text(slide, Inches(0.6), Inches(2.5), Inches(12.4), Inches(2),
             data.get("headline", "Obrigado."),
             size=72, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(4.5), Inches(12.4), Inches(0.6),
             data.get("sub", ""), size=18, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(11), Inches(6.7), Inches(2), Inches(0.3),
             data.get("brand", "Singular Group"), size=10, italic=True, color=MUTED, align=PP_ALIGN.RIGHT)


def add_slide_number(slide, idx, total):
    add_text(slide, Inches(11.3), Inches(0.4), Inches(2), Inches(0.3),
             f"{idx+1:02d} / {total:02d}", size=10, color=MUTED, align=PP_ALIGN.RIGHT, letter_spacing=2)


RENDERERS = {
    "cover": lambda s, d, t, i: render_cover(s, d, t),
    "summary-cards": render_summary_cards,
    "section-detail": render_section_detail,
    "compare-table": render_compare_table,
    "text-section": render_text_section,
    "recommendation": render_recommendation,
    "next-steps": render_next_steps,
    "closing": render_closing,
}


def build(content: dict, output_path: Path):
    slides_data = content.get("slides", [])
    if not slides_data:
        raise SystemExit("content.json: lista de slides vazia.")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    total = len(slides_data)
    for i, sd in enumerate(slides_data):
        stype = sd.get("type")
        renderer = RENDERERS.get(stype)
        if not renderer:
            print(f"Aviso: tipo de slide desconhecido '{stype}' (slide {i+1}). Pulado.", file=sys.stderr)
            continue
        slide = prs.slides.add_slide(blank_layout)
        renderer(slide, sd, total, i)

    prs.save(str(output_path))


def main():
    if len(sys.argv) < 3:
        print("Uso: python build_pptx.py content.json saida.pptx", file=sys.stderr)
        sys.exit(1)
    content_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2])
    if not content_path.exists():
        print(f"Erro: nao encontrei {content_path}", file=sys.stderr)
        sys.exit(2)
    with content_path.open("r", encoding="utf-8") as f:
        content = json.load(f)
    build(content, output_path)
    size_kb = output_path.stat().st_size / 1024
    n_slides = len(content.get("slides", []))
    print(f"Gerado: {output_path} ({size_kb:.1f}KB, {n_slides} slides)")


if __name__ == "__main__":
    main()
